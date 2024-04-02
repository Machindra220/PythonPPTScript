# importing 
from pptx import Presentation 
from pptx.util import Inches 

# create a Presentation object 
ppt = Presentation() 

# Adding a blank slide in out ppt 
slide = ppt.slides.add_slide(ppt.slide_layouts[0]) 

""" Ref for slide types: 
0 -> title and subtitle 
1 -> title and content 
2 -> section header 
3 -> two content 
4 -> Comparison 
5 -> Title only 
6 -> Blank 
7 -> Content with caption 
8 -> Pic with caption 
"""


# slide i.e. first page of slide 
slide.shapes.title.text = " Problem Tickets"
slide.placeholders[1].text = " NEWLY CREATED / UPDATED"

# Adjusting the width ! Where x=leftMargin, y=topMargin, cx=WidthOfTable, cy=HeightofTable
x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4) 

# Adding tables 
shape = slide.shapes.add_table(4, 3, x, y, cx, cy) 
table = shape.table
cell = table.cell(0, 0)
cell.text = 'Problem Ticket #'
cell = table.cell(0, 1)
cell.text = 'Configuration Item'
cell = table.cell(0, 2)
cell.text = 'Comments'

# Saving the file 
ppt.save("./Outputs/Tabel_Tutorial9c.pptx") 

print("done")
