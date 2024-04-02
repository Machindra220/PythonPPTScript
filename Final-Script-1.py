# import required things 
from pptx import Presentation 
from pptx.util import Inches, Pt, Cm 

# Creating Object 
ppt = Presentation() 
# set width and height to 16 and 9 inches.
ppt.slide_width = Inches(16)
ppt.slide_height = Inches(9)

# To create blank slide layout 
# We have to use 6 as an argument of slide_layouts (Blank Slide)
blank_slide_layout = ppt.slide_layouts[6] 

# =======Slide 1 Start ============
# Attaching slide obj to slide 1
slide = ppt.slides.add_slide(blank_slide_layout) 

# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(7), Cm(8), Cm(25), Cm(2.5) 

# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 

# creating textFrames 
tf = txBox.text_frame 
#tf.text = "This is text inside a textbox"
# adding Paragraphs 
p = tf.add_paragraph() 
p.font.bold = True
p.font.size = Pt(40)
#p.font.italic = True
p.text = "Platform Operations Weekly Report"

# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(12), Cm(12), Cm(12), Cm(2.5) 

# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.font.size = Pt(28)
p.text = "12/11/2023 - 12/17/2023"

# =======Slide 1 Complete==========
# =======Slide 2 Start ============
slide = ppt.slides.add_slide(blank_slide_layout) 

# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)

# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 

# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PRESENTATION SCHEDULE"
p.font.size = Pt(38)
p.font.bold = True
#Slide 2 Title complete
# adding Paragraphs 
#p = tf.add_paragraph() 
# adding text 
#p.text = "This is a second paragraph that's bold and italic"
# font 
#p.font.bold = True
#p.font.italic = True
#p = tf.add_paragraph() 
#p.text = "This is a third paragraph that's big "
#p.font.size = Pt(40) 
#Table Start
x, y, cx, cy = Cm(1.5), Cm(3), Cm(32), Cm(12) # x=left, y=height, cx=widthOfCol, cy=HeightOfRow
shape = slide.shapes.add_table(6, 3, x, y, cx, cy) #4 = rows, 5=columns
if shape.has_table:
    table = shape.table
#1st Row
cell = table.cell(0, 0)  # Access the cell at row 1, column 1 (0-based index)
cell.text = "Date"
cell = table.cell(0, 1)  # Access the cell at row 1, column 2 (0-based index)
cell.text = "Presenter"
cell = table.cell(0, 2)  # Access the cell at row 1, column 3 (0-based index)
cell.text = "Status"
#2nd Row
cell = table.cell(1, 0)  # Access the cell at row 2, column 1 (0-based index)
cell.text = "18/12/2023"
cell = table.cell(1, 1)  # Access the cell at row 2, column 2 (0-based index)
cell.text = "Person-1 / Person-2"
cell = table.cell(1, 2)  # Access the cell at row 2, column 3 (0-based index)
cell.text = "Done"
#3rd Row
cell = table.cell(2, 0)  # Access the cell at row 3, column 1 (0-based index)
cell.text = "18/12/2024"
cell = table.cell(2, 1)  # Access the cell at row 3, column 2 (0-based index)
cell.text = "Person-1 / Person-3"
cell = table.cell(2, 2)  # Access the cell at row 3, column 3 (0-based index)
cell.text = "Done"
#Table Complete
#Slide 2 complete
#Slide 3 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "Summary - "
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(8), Cm(0.7), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "COUNT/PRIORITY/STATUS (ENVIRONMENT-WISE) "
p.font.size = Pt(20)
#Slide 3 Title complete
#Slide 3, 2nd small centered Title
left, top, width, height = Cm(16), Cm(1.9), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "CURRENT WEEK "
p.font.size = Pt(16)
# Giving Image-1 path 
img_path = './Images/bg_bg1.jpeg'
# For Image-1 margins 
left, top, width, height = Cm(2), Cm(3.5), Cm(36), Cm(9)
pic = slide.shapes.add_picture(img_path, left, top, width, height)

#Slide 3 Middle Title
left, top, width, height = Cm(16), Cm(12), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PAST WEEK "
p.font.size = Pt(16)

# Giving Image-2 path 
img_path = './Images/bg_bg.jpeg'
# For Image-2 margins 
left, top, width, height = Cm(2), Cm(14), Cm(36), Cm(8)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
#Slide 3 Complete

#Slide 4 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "AUTO-RESOLVED METRICS - "
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(17), Cm(0.7), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "COUNT / CATEGORY "
p.font.size = Pt(20)
#Slide 4 Title complete
#Slide 4, 2nd small centered Title
left, top, width, height = Cm(16), Cm(1.9), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "CURRENT WEEK "
p.font.size = Pt(16)
# Giving Image-1 path 
img_path = './Images/bg_bg1.jpeg'
# For Image-1 margins 
left, top, width, height = Cm(2), Cm(3.5), Cm(36), Cm(9)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
#Slide 4 Middle Title
left, top, width, height = Cm(16), Cm(12), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PAST WEEK "
p.font.size = Pt(16)
# Giving Image-2 path 
img_path = './Images/bg_bg.jpeg'
# For Image-2 margins 
left, top, width, height = Cm(2), Cm(14), Cm(36), Cm(8)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
#Slide 4 Complete
#Slide 5 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "AUTO-RESOLVED METRICS - "
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(17), Cm(0.7), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "COUNT / CATEGORY "
p.font.size = Pt(20)
#Slide 5 Title complete
#Slide 5, 2nd small centered Title
left, top, width, height = Cm(16), Cm(1.9), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "CURRENT WEEK "
p.font.size = Pt(16)
# Slide 5, Image-1 path 
img_path = './Images/bg_bg1.jpeg'
# Slide 5, Image-1 margins 
left, top, width, height = Cm(2), Cm(3.5), Cm(36), Cm(4)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
# Slide 5, Image-2 path 
img_path = './Images/bg_bg1.jpeg'
# Slide 5, Image-2 margins 
left, top, width, height = Cm(2), Cm(8.1), Cm(36), Cm(4.5)
pic = slide.shapes.add_picture(img_path, left, top, width, height)

#Slide 5 Middle Title
left, top, width, height = Cm(16), Cm(12), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PAST WEEK "
p.font.size = Pt(16)
# Slide 5, Image-3 path 
img_path = './Images/bg_bg.jpeg'
# Slide 5, Image-3 margins 
left, top, width, height = Cm(2), Cm(13.4), Cm(36), Cm(4)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
# Slide 5, Image-4 path 
img_path = './Images/bg_bg.jpeg'
# Slide 5, Image-4 margins 
left, top, width, height = Cm(2), Cm(17.7), Cm(36), Cm(4.5)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
#Slide 5 Complete
#Slide 6 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "MANUALLY-RESOLVED METRICS - "
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(20), Cm(0.7), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "COUNT / CATEGORY "
p.font.size = Pt(20)
#Slide 6 Title complete
#Slide 6, 2nd small centered Title
left, top, width, height = Cm(16), Cm(1.9), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "CURRENT WEEK "
p.font.size = Pt(16)

#Slide 6 Middle Title
left, top, width, height = Cm(16), Cm(12), Cm(10), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PAST WEEK "
p.font.size = Pt(16)

#Slide 6 Complete
#Slide 7 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PROBLEMS TICKETS"
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(1.3), Cm(2), Cm(20), Cm(1.4)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "NEWLY CREATED / UPDATED "
p.font.size = Pt(24)
#Slide 7 Title complete

#Slide 7 Complete
#Slide 8 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "RUNBOOK METRICS"
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(1.3), Cm(2), Cm(20), Cm(1.2)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "NEWLY CREATED / UPDATED "
p.font.size = Pt(24)
#Slide 8 Title complete

#Slide 8 Complete
#Slide 9 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1.4), Cm(0.2), Cm(20), Cm(1.8)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "KNOWLEDGE ARTICLES METRICS "
p.font.size = Pt(38)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(1.3), Cm(2), Cm(20), Cm(1.2)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "NEWLY CREATED / UPDATED "
p.font.size = Pt(24)
#Slide 9 Title complete

#Slide 9 Complete
#Slide 10 Start
slide = ppt.slides.add_slide(blank_slide_layout) 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(12), Cm(8), Cm(16), Cm(4)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "THANK YOU "
p.font.size = Pt(88)
p.font.bold = True
#Small Font Title
left, top, width, height = Cm(17), Cm(12), Cm(8), Cm(1.5)
# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "TEAM : Operations "
p.font.size = Pt(24)
#Slide 10 Title complete

#Slide 10 Complete

# save file 
ppt.save('Presentation_Output.pptx') 

print("done") 
