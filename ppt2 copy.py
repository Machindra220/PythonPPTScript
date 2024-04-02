# import required things 
from pptx import Presentation 
from pptx.util import Inches, Pt, Cm 

# Creating Object 
ppt = Presentation() 
# set width and height to 16 and 9 inches.
ppt.slide_width = Inches(16)
ppt.slide_height = Inches(9)

# To create blank slide layout 
# We have to use 6 as an argument 
# of slide_layouts 
blank_slide_layout = ppt.slide_layouts[6] 

# Attaching slide obj to slide 1
slide = ppt.slides.add_slide(blank_slide_layout) 

# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(7), Cm(8), Cm(25), Cm(2.5) 

# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 

# creating textFrames 
tf = txBox.text_frame 
#tf.text = "This is text inside a textbox"
# adding Paragraphs 
p = tf.add_paragraph() 
p.font.bold = True
#p.font.italic = True
#p = tf.add_paragraph() 
p.text = "Platform Operations: L2 Incident Report"
p.font.size = Pt(40)

# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(12), Cm(12), Cm(12), Cm(2.5) 

# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 
# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "12/11/2023 - 12/17/2023"
p.font.size = Pt(28)

# Slide 2
slide = ppt.slides.add_slide(blank_slide_layout) 

# For adjusting the Margins in inches 
# For adjusting the Margins in CM/inches 
left, top, width, height = Cm(1), Cm(0.2), Cm(20), Cm(1.8)

# creating textBox 
txBox = slide.shapes.add_textbox(left, top, width, height) 

# creating textFrames 
tf = txBox.text_frame 
p = tf.add_paragraph()
p.text = "PRESENTATION SCHEDULE"
p.font.size = Pt(38)
p.font.bold = True

# adding Paragraphs 
p = tf.add_paragraph() 
# adding text 
p.text = "This is a second paragraph that's bold and italic"
# font 
p.font.bold = True
p.font.italic = True
p = tf.add_paragraph() 
p.text = "This is a third paragraph that's big "
p.font.size = Pt(40) 

# save file 
ppt.save('test_3b.pptx') 

print("done") 
