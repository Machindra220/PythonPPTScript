from pptx import Presentation 
from pptx.util import Inches, Pt, Cm 

# Giving Image path 
img_path = './Images/bg_bg.jpeg'

# Creating an Presentation object 
ppt = Presentation() 
# set width and height to 16 and 9 inches.
ppt.slide_width = Inches(16)
ppt.slide_height = Inches(9)

# Selecting blank slide 
blank_slide_layout = ppt.slide_layouts[6] 

# Attaching slide to ppt 
slide = ppt.slides.add_slide(blank_slide_layout) 

# For margins 
left, top, width, height = Cm(4), Cm(2), Cm(30), Cm(28)


# adding images 
#pic = slide.shapes.add_picture(img_path,left, top) 

#left = Inches(0) 
#height = Inches(0) 

pic = slide.shapes.add_picture(img_path, left, 
							top, width, height) 
# save file 
ppt.save('./Outputs/test_45.pptx') 

print("Done")
