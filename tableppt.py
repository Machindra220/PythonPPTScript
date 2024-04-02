# importing 
from pptx import Presentation 
from pptx.util import Inches 

# create a Presentation object 
ppt = Presentation() 

# Adding a blank slide in out ppt 
slide = ppt.slides.add_slide(ppt.slide_layouts[6]) 

# Adjusting the width ! Where x=leftMargin, y=topMargin, cx=WidthOfTable, cy=HeightofTable
x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(4) 

# Adding tables 
shape = slide.shapes.add_table(4, 4, x, y, cx, cy) 

# Saving the file 
ppt.save("./Outputs/Tabel_Tutorial8.pptx") 

print("done")
